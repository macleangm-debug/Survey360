"""DataPulse - Data Analysis Module
Comprehensive analytics engine for survey data analysis.

This module provides:
- Response browsing and filtering
- Dataset snapshots for reproducible analysis
- Statistical analysis (descriptive, inferential, regression)
- AI-powered natural language analysis
- Export capabilities
"""

from fastapi import APIRouter, HTTPException, Request, Query, BackgroundTasks
from pydantic import BaseModel, Field
from typing import List, Optional, Dict, Any, Union
from datetime import datetime, timezone, timedelta
from enum import Enum
import pandas as pd
import numpy as np
import json
import hashlib
import io

router = APIRouter(prefix="/analysis", tags=["Data Analysis"])


# ============ Enums ============

class VariableType(str, Enum):
    STRING = "string"
    NUMERIC = "numeric"
    DATE = "date"
    CATEGORICAL = "categorical"


class MeasurementLevel(str, Enum):
    NOMINAL = "nominal"
    ORDINAL = "ordinal"
    SCALE = "scale"


class SnapshotStatus(str, Enum):
    CREATING = "creating"
    READY = "ready"
    FAILED = "failed"
    ARCHIVED = "archived"


# ============ Models ============

class ResponseFilter(BaseModel):
    form_id: str
    date_from: Optional[str] = None
    date_to: Optional[str] = None
    status: Optional[List[str]] = None
    language: Optional[str] = None
    field_filters: Optional[Dict[str, Any]] = None  # {"field_id": {"operator": "eq", "value": "x"}}
    page: int = 1
    page_size: int = 50


class SnapshotConfig(BaseModel):
    form_id: str
    org_id: str
    name: str
    description: Optional[str] = None
    include_statuses: List[str] = ["approved"]
    include_metadata: bool = True
    include_paradata: bool = False
    anonymize_pii: bool = False


class QuickStatsRequest(BaseModel):
    snapshot_id: Optional[str] = None
    form_id: Optional[str] = None
    org_id: str
    variables: List[str]  # Field IDs to analyze
    group_by: Optional[str] = None  # Optional grouping variable


class CrosstabRequest(BaseModel):
    snapshot_id: Optional[str] = None
    form_id: Optional[str] = None
    org_id: str
    row_var: str
    col_var: str
    weight_var: Optional[str] = None


class ExportRequest(BaseModel):
    snapshot_id: Optional[str] = None
    form_id: Optional[str] = None
    org_id: str
    format: str = "csv"  # csv, xlsx, spss, stata, r
    include_labels: bool = True
    include_codebook: bool = False
    variables: Optional[List[str]] = None  # None = all


# ============ Response Browsing ============

@router.post("/responses/browse")
async def browse_responses(
    request: Request,
    filters: ResponseFilter
):
    """Browse and filter individual responses"""
    db = request.app.state.db
    
    # Build query
    query = {"form_id": filters.form_id}
    
    if filters.date_from:
        query["submitted_at"] = {"$gte": datetime.fromisoformat(filters.date_from)}
    if filters.date_to:
        if "submitted_at" in query:
            query["submitted_at"]["$lte"] = datetime.fromisoformat(filters.date_to)
        else:
            query["submitted_at"] = {"$lte": datetime.fromisoformat(filters.date_to)}
    
    if filters.status:
        query["status"] = {"$in": filters.status}
    
    if filters.language:
        query["language"] = filters.language
    
    # Field-level filters
    if filters.field_filters:
        for field_id, condition in filters.field_filters.items():
            op = condition.get("operator", "eq")
            val = condition.get("value")
            
            if op == "eq":
                query[f"data.{field_id}"] = val
            elif op == "ne":
                query[f"data.{field_id}"] = {"$ne": val}
            elif op == "gt":
                query[f"data.{field_id}"] = {"$gt": val}
            elif op == "lt":
                query[f"data.{field_id}"] = {"$lt": val}
            elif op == "contains":
                query[f"data.{field_id}"] = {"$regex": val, "$options": "i"}
            elif op == "in":
                query[f"data.{field_id}"] = {"$in": val}
    
    # Execute query with pagination
    skip = (filters.page - 1) * filters.page_size
    
    total = await db.submissions.count_documents(query)
    submissions = await db.submissions.find(query).skip(skip).limit(filters.page_size).to_list(filters.page_size)
    
    # Format response
    for s in submissions:
        s["_id"] = str(s.get("_id", ""))
        if s.get("submitted_at"):
            s["submitted_at"] = s["submitted_at"].isoformat()
        if s.get("created_at"):
            s["created_at"] = s["created_at"].isoformat()
    
    return {
        "total": total,
        "page": filters.page,
        "page_size": filters.page_size,
        "total_pages": (total + filters.page_size - 1) // filters.page_size,
        "responses": submissions
    }


@router.get("/responses/{submission_id}")
async def get_response_detail(
    request: Request,
    submission_id: str
):
    """Get detailed view of a single response with metadata"""
    db = request.app.state.db
    
    submission = await db.submissions.find_one({"id": submission_id})
    if not submission:
        raise HTTPException(status_code=404, detail="Response not found")
    
    submission["_id"] = str(submission.get("_id", ""))
    
    # Get paradata if available
    paradata = await db.paradata.find_one({"submission_id": submission_id})
    if paradata:
        paradata["_id"] = str(paradata.get("_id", ""))
        submission["paradata"] = paradata
    
    # Get revision history
    revisions = await db.revisions.find({"submission_id": submission_id}).to_list(50)
    for r in revisions:
        r["_id"] = str(r.get("_id", ""))
    submission["revisions"] = revisions
    
    return submission


@router.get("/responses/{form_id}/summary")
async def get_response_summary(
    request: Request,
    form_id: str,
    org_id: str
):
    """Get quick summary statistics for a form's responses"""
    db = request.app.state.db
    
    # Aggregation pipeline
    pipeline = [
        {"$match": {"form_id": form_id}},
        {"$group": {
            "_id": None,
            "total": {"$sum": 1},
            "approved": {"$sum": {"$cond": [{"$eq": ["$status", "approved"]}, 1, 0]}},
            "pending": {"$sum": {"$cond": [{"$eq": ["$status", "submitted"]}, 1, 0]}},
            "rejected": {"$sum": {"$cond": [{"$eq": ["$status", "rejected"]}, 1, 0]}},
            "draft": {"$sum": {"$cond": [{"$eq": ["$status", "draft"]}, 1, 0]}},
            "first_submission": {"$min": "$submitted_at"},
            "last_submission": {"$max": "$submitted_at"}
        }}
    ]
    
    result = await db.submissions.aggregate(pipeline).to_list(1)
    
    if not result:
        return {
            "total": 0,
            "by_status": {},
            "date_range": None
        }
    
    stats = result[0]
    return {
        "total": stats["total"],
        "by_status": {
            "approved": stats["approved"],
            "pending": stats["pending"],
            "rejected": stats["rejected"],
            "draft": stats["draft"]
        },
        "date_range": {
            "first": stats["first_submission"].isoformat() if stats["first_submission"] else None,
            "last": stats["last_submission"].isoformat() if stats["last_submission"] else None
        }
    }


# ============ Dataset Snapshots ============

@router.post("/snapshots/create")
async def create_snapshot(
    request: Request,
    config: SnapshotConfig,
    background_tasks: BackgroundTasks
):
    """Create an immutable dataset snapshot for analysis"""
    db = request.app.state.db
    
    snapshot_id = f"snap_{config.form_id}_{int(datetime.now(timezone.utc).timestamp())}"
    
    # Create snapshot record
    snapshot = {
        "id": snapshot_id,
        "form_id": config.form_id,
        "org_id": config.org_id,
        "name": config.name,
        "description": config.description,
        "config": config.dict(),
        "status": SnapshotStatus.CREATING,
        "row_count": 0,
        "created_at": datetime.now(timezone.utc),
        "created_by": None,  # TODO: Add user
        "schema_hash": None,
        "data_hash": None
    }
    
    await db.snapshots.insert_one(snapshot)
    
    # Process in background
    background_tasks.add_task(
        process_snapshot_creation,
        db, snapshot_id, config
    )
    
    return {
        "snapshot_id": snapshot_id,
        "status": "creating",
        "message": "Snapshot creation started"
    }


async def process_snapshot_creation(db, snapshot_id: str, config: SnapshotConfig):
    """Background task to create snapshot data"""
    try:
        # Get form schema
        form = await db.forms.find_one({"id": config.form_id})
        if not form:
            await db.snapshots.update_one(
                {"id": snapshot_id},
                {"$set": {"status": SnapshotStatus.FAILED, "error": "Form not found"}}
            )
            return
        
        # Query submissions
        query = {
            "form_id": config.form_id,
            "status": {"$in": config.include_statuses}
        }
        
        submissions = await db.submissions.find(query).to_list(None)
        
        # Build snapshot data
        snapshot_data = []
        for sub in submissions:
            row = {
                "_submission_id": sub["id"],
                "_submitted_at": sub.get("submitted_at", "").isoformat() if sub.get("submitted_at") else None,
                "_status": sub.get("status"),
            }
            
            if config.include_metadata:
                row["_language"] = sub.get("language")
                row["_device"] = sub.get("device_type")
                row["_duration_seconds"] = sub.get("duration_seconds")
            
            # Add response data
            data = sub.get("data", {})
            for key, value in data.items():
                if config.anonymize_pii:
                    # Check if field is PII (simplified check)
                    field = next((f for f in form.get("fields", []) if f["id"] == key), None)
                    if field and field.get("pii"):
                        row[key] = "[REDACTED]"
                        continue
                row[key] = value
            
            snapshot_data.append(row)
        
        # Calculate hashes for reproducibility
        schema_str = json.dumps(form.get("fields", []), sort_keys=True)
        schema_hash = hashlib.sha256(schema_str.encode()).hexdigest()[:16]
        
        data_str = json.dumps(snapshot_data, sort_keys=True, default=str)
        data_hash = hashlib.sha256(data_str.encode()).hexdigest()[:16]
        
        # Store snapshot data
        await db.snapshot_data.insert_one({
            "snapshot_id": snapshot_id,
            "data": snapshot_data,
            "schema": form.get("fields", [])
        })
        
        # Update snapshot status
        await db.snapshots.update_one(
            {"id": snapshot_id},
            {"$set": {
                "status": SnapshotStatus.READY,
                "row_count": len(snapshot_data),
                "schema_hash": schema_hash,
                "data_hash": data_hash,
                "completed_at": datetime.now(timezone.utc)
            }}
        )
        
    except Exception as e:
        await db.snapshots.update_one(
            {"id": snapshot_id},
            {"$set": {"status": SnapshotStatus.FAILED, "error": str(e)}}
        )


@router.get("/snapshots/{org_id}")
async def list_snapshots(
    request: Request,
    org_id: str,
    form_id: Optional[str] = None
):
    """List all snapshots for an organization"""
    db = request.app.state.db
    
    query = {"org_id": org_id}
    if form_id:
        query["form_id"] = form_id
    
    snapshots = await db.snapshots.find(query).sort("created_at", -1).to_list(100)
    
    for s in snapshots:
        s["_id"] = str(s.get("_id", ""))
        if s.get("created_at"):
            s["created_at"] = s["created_at"].isoformat()
        if s.get("completed_at"):
            s["completed_at"] = s["completed_at"].isoformat()
    
    return snapshots


@router.get("/snapshots/detail/{snapshot_id}")
async def get_snapshot_detail(
    request: Request,
    snapshot_id: str
):
    """Get snapshot metadata and schema"""
    db = request.app.state.db
    
    snapshot = await db.snapshots.find_one({"id": snapshot_id})
    if not snapshot:
        raise HTTPException(status_code=404, detail="Snapshot not found")
    
    snapshot["_id"] = str(snapshot.get("_id", ""))
    
    # Get schema
    snapshot_data = await db.snapshot_data.find_one({"snapshot_id": snapshot_id})
    if snapshot_data:
        snapshot["schema"] = snapshot_data.get("schema", [])
        snapshot["preview"] = snapshot_data.get("data", [])[:5]  # First 5 rows
    
    return snapshot


# ============ Quick Statistics ============

@router.post("/stats/quick")
async def get_quick_stats(
    request: Request,
    req: QuickStatsRequest
):
    """Get quick statistics for selected variables"""
    db = request.app.state.db
    
    # Get data
    if req.snapshot_id:
        snapshot_data = await db.snapshot_data.find_one({"snapshot_id": req.snapshot_id})
        if not snapshot_data:
            raise HTTPException(status_code=404, detail="Snapshot not found")
        data = snapshot_data.get("data", [])
        schema = snapshot_data.get("schema", [])
    else:
        # Use live data
        submissions = await db.submissions.find({
            "form_id": req.form_id,
            "status": "approved"
        }).to_list(None)
        data = [s.get("data", {}) for s in submissions]
        form = await db.forms.find_one({"id": req.form_id})
        schema = form.get("fields", []) if form else []
    
    if not data:
        return {"variables": [], "total_n": 0}
    
    # Convert to DataFrame
    df = pd.DataFrame(data)
    
    # Calculate stats for each variable
    results = []
    for var_id in req.variables:
        if var_id not in df.columns:
            continue
        
        field_schema = next((f for f in schema if f.get("id") == var_id), {})
        field_type = field_schema.get("type", "text")
        
        series = df[var_id].dropna()
        
        if field_type in ["number", "integer", "decimal"]:
            # Numeric statistics
            numeric_series = pd.to_numeric(series, errors='coerce').dropna()
            stats = {
                "variable": var_id,
                "label": field_schema.get("label", var_id),
                "type": "numeric",
                "n": len(numeric_series),
                "missing": len(df) - len(numeric_series),
                "mean": float(numeric_series.mean()) if len(numeric_series) > 0 else None,
                "median": float(numeric_series.median()) if len(numeric_series) > 0 else None,
                "std": float(numeric_series.std()) if len(numeric_series) > 1 else None,
                "min": float(numeric_series.min()) if len(numeric_series) > 0 else None,
                "max": float(numeric_series.max()) if len(numeric_series) > 0 else None,
                "percentiles": {
                    "25": float(numeric_series.quantile(0.25)) if len(numeric_series) > 0 else None,
                    "50": float(numeric_series.quantile(0.50)) if len(numeric_series) > 0 else None,
                    "75": float(numeric_series.quantile(0.75)) if len(numeric_series) > 0 else None
                }
            }
        else:
            # Categorical statistics
            value_counts = series.value_counts()
            total = len(series)
            
            frequencies = []
            for val, count in value_counts.items():
                frequencies.append({
                    "value": str(val),
                    "count": int(count),
                    "percent": round(count / total * 100, 1) if total > 0 else 0
                })
            
            stats = {
                "variable": var_id,
                "label": field_schema.get("label", var_id),
                "type": "categorical",
                "n": len(series),
                "missing": len(df) - len(series),
                "unique_values": len(value_counts),
                "frequencies": frequencies[:20],  # Top 20
                "mode": str(value_counts.index[0]) if len(value_counts) > 0 else None
            }
        
        results.append(stats)
    
    return {
        "total_n": len(df),
        "variables": results
    }


@router.post("/stats/crosstab")
async def get_crosstab(
    request: Request,
    req: CrosstabRequest
):
    """Generate a cross-tabulation with chi-square test"""
    db = request.app.state.db
    
    # Get data
    if req.snapshot_id:
        snapshot_data = await db.snapshot_data.find_one({"snapshot_id": req.snapshot_id})
        if not snapshot_data:
            raise HTTPException(status_code=404, detail="Snapshot not found")
        data = snapshot_data.get("data", [])
    else:
        submissions = await db.submissions.find({
            "form_id": req.form_id,
            "status": "approved"
        }).to_list(None)
        data = [s.get("data", {}) for s in submissions]
    
    if not data:
        return {"error": "No data available"}
    
    df = pd.DataFrame(data)
    
    if req.row_var not in df.columns or req.col_var not in df.columns:
        raise HTTPException(status_code=400, detail="Variable not found in data")
    
    # Create crosstab
    crosstab = pd.crosstab(
        df[req.row_var],
        df[req.col_var],
        margins=True,
        margins_name="Total"
    )
    
    # Chi-square test
    from scipy import stats as scipy_stats
    
    # Remove margins for chi-square
    ct_no_margins = crosstab.iloc[:-1, :-1]
    
    chi2_result = None
    if ct_no_margins.shape[0] > 1 and ct_no_margins.shape[1] > 1:
        try:
            chi2, p_value, dof, expected = scipy_stats.chi2_contingency(ct_no_margins)
            chi2_result = {
                "chi_square": round(float(chi2), 4),
                "p_value": round(float(p_value), 4),
                "degrees_of_freedom": int(dof),
                "significant": bool(p_value < 0.05)
            }
        except Exception:
            pass
    
    # Format output
    result = {
        "row_variable": req.row_var,
        "col_variable": req.col_var,
        "row_labels": [str(x) for x in crosstab.index.tolist()],
        "col_labels": [str(x) for x in crosstab.columns.tolist()],
        "counts": [[int(c) for c in row] for row in crosstab.values.tolist()],
        "chi_square_test": chi2_result
    }
    
    # Add percentages
    row_pcts = (crosstab.div(crosstab.iloc[:, -1], axis=0) * 100).round(1)
    result["row_percentages"] = [[float(c) for c in row] for row in row_pcts.values.tolist()]
    
    col_pcts = (crosstab.div(crosstab.iloc[-1, :], axis=1) * 100).round(1)
    result["col_percentages"] = [[float(c) for c in row] for row in col_pcts.values.tolist()]
    
    return result


# ============ Time Series / Trends ============

@router.get("/stats/trends/{form_id}")
async def get_submission_trends(
    request: Request,
    form_id: str,
    days: int = 30,
    group_by: Optional[str] = None
):
    """Get submission trends over time"""
    db = request.app.state.db
    
    start_date = datetime.now(timezone.utc) - timedelta(days=days)
    
    pipeline = [
        {"$match": {
            "form_id": form_id,
            "submitted_at": {"$gte": start_date}
        }},
        {"$group": {
            "_id": {
                "date": {"$dateToString": {"format": "%Y-%m-%d", "date": "$submitted_at"}}
            },
            "count": {"$sum": 1},
            "approved": {"$sum": {"$cond": [{"$eq": ["$status", "approved"]}, 1, 0]}},
            "pending": {"$sum": {"$cond": [{"$eq": ["$status", "submitted"]}, 1, 0]}}
        }},
        {"$sort": {"_id.date": 1}}
    ]
    
    results = await db.submissions.aggregate(pipeline).to_list(None)
    
    trends = []
    for r in results:
        trends.append({
            "date": r["_id"]["date"],
            "count": r["count"],
            "approved": r["approved"],
            "pending": r["pending"]
        })
    
    return {"trends": trends, "period_days": days}


# ============ Missing Data Imputation ============

class ImputationMethod(str, Enum):
    MEAN = "mean"
    MEDIAN = "median"
    MODE = "mode"
    CONSTANT = "constant"
    FORWARD_FILL = "ffill"
    BACKWARD_FILL = "bfill"
    INTERPOLATE = "interpolate"
    DROP = "drop"


class ImputationRequest(BaseModel):
    org_id: str
    form_id: Optional[str] = None
    snapshot_id: Optional[str] = None
    variables: List[str]
    method: ImputationMethod = ImputationMethod.MEAN
    constant_value: Optional[Any] = None
    group_by: Optional[str] = None  # For group-wise imputation
    create_snapshot: bool = False  # Create new snapshot with imputed data


@router.post("/imputation/preview")
async def preview_imputation(
    request: Request,
    req: ImputationRequest
):
    """Preview missing data imputation without applying changes"""
    db = request.app.state.db
    
    # Get data
    if req.snapshot_id:
        snapshot_data = await db.snapshot_data.find_one({"snapshot_id": req.snapshot_id})
        if not snapshot_data:
            raise HTTPException(status_code=404, detail="Snapshot not found")
        data = snapshot_data.get("data", [])
    else:
        submissions = await db.submissions.find({
            "form_id": req.form_id,
            "status": {"$in": ["approved", "submitted"]}
        }).to_list(None)
        data = [s.get("data", {}) for s in submissions]
    
    if not data:
        return {"error": "No data available"}
    
    df = pd.DataFrame(data)
    
    # Filter to requested variables
    valid_vars = [v for v in req.variables if v in df.columns]
    if not valid_vars:
        return {"error": "No valid variables found"}
    
    # Calculate missing data statistics before imputation
    missing_before = {}
    for var in valid_vars:
        missing_count = df[var].isna().sum() + (df[var] == '').sum()
        missing_before[var] = {
            "count": int(missing_count),
            "percent": round(float(missing_count / len(df) * 100), 2)
        }
    
    # Apply imputation to copy of data
    df_imputed = df.copy()
    imputation_details = {}
    
    for var in valid_vars:
        # Convert to numeric if possible for numeric methods
        series = df_imputed[var].replace('', np.nan)
        try:
            numeric_series = pd.to_numeric(series, errors='coerce')
            is_numeric = numeric_series.notna().sum() > 0
        except Exception:
            is_numeric = False
            numeric_series = series
        
        imputed_value = None
        
        if req.method == ImputationMethod.MEAN:
            if is_numeric:
                if req.group_by and req.group_by in df_imputed.columns:
                    df_imputed[var] = df_imputed.groupby(req.group_by)[var].transform(
                        lambda x: pd.to_numeric(x, errors='coerce').fillna(pd.to_numeric(x, errors='coerce').mean())
                    )
                    imputed_value = "Group means"
                else:
                    imputed_value = float(numeric_series.mean())
                    df_imputed[var] = numeric_series.fillna(imputed_value)
            else:
                imputed_value = "N/A (not numeric)"
                
        elif req.method == ImputationMethod.MEDIAN:
            if is_numeric:
                if req.group_by and req.group_by in df_imputed.columns:
                    df_imputed[var] = df_imputed.groupby(req.group_by)[var].transform(
                        lambda x: pd.to_numeric(x, errors='coerce').fillna(pd.to_numeric(x, errors='coerce').median())
                    )
                    imputed_value = "Group medians"
                else:
                    imputed_value = float(numeric_series.median())
                    df_imputed[var] = numeric_series.fillna(imputed_value)
            else:
                imputed_value = "N/A (not numeric)"
                
        elif req.method == ImputationMethod.MODE:
            mode_val = series.mode()
            if len(mode_val) > 0:
                imputed_value = str(mode_val.iloc[0])
                df_imputed[var] = series.fillna(mode_val.iloc[0])
            else:
                imputed_value = "N/A (no mode)"
                
        elif req.method == ImputationMethod.CONSTANT:
            imputed_value = req.constant_value
            df_imputed[var] = series.fillna(req.constant_value)
            
        elif req.method == ImputationMethod.FORWARD_FILL:
            df_imputed[var] = series.ffill()
            imputed_value = "Forward fill"
            
        elif req.method == ImputationMethod.BACKWARD_FILL:
            df_imputed[var] = series.bfill()
            imputed_value = "Backward fill"
            
        elif req.method == ImputationMethod.INTERPOLATE:
            if is_numeric:
                df_imputed[var] = numeric_series.interpolate(method='linear')
                imputed_value = "Linear interpolation"
            else:
                imputed_value = "N/A (not numeric)"
                
        elif req.method == ImputationMethod.DROP:
            imputed_value = "Rows dropped"
        
        imputation_details[var] = {
            "method": req.method.value,
            "imputed_value": imputed_value,
            "is_numeric": is_numeric
        }
    
    # Handle drop method
    if req.method == ImputationMethod.DROP:
        original_count = len(df_imputed)
        df_imputed = df_imputed.dropna(subset=valid_vars)
        rows_dropped = original_count - len(df_imputed)
    else:
        rows_dropped = 0
    
    # Calculate missing data statistics after imputation
    missing_after = {}
    for var in valid_vars:
        if var in df_imputed.columns:
            missing_count = df_imputed[var].isna().sum()
            missing_after[var] = {
                "count": int(missing_count),
                "percent": round(float(missing_count / len(df_imputed) * 100), 2) if len(df_imputed) > 0 else 0
            }
    
    # Preview sample of changes
    sample_changes = []
    for idx in df.index[:20]:
        if idx in df_imputed.index:
            for var in valid_vars:
                original = df.loc[idx, var] if var in df.columns else None
                imputed = df_imputed.loc[idx, var] if var in df_imputed.columns else None
                if pd.isna(original) or original == '':
                    sample_changes.append({
                        "row": int(idx),
                        "variable": var,
                        "original": None,
                        "imputed": imputed if not pd.isna(imputed) else None
                    })
    
    return {
        "n_original": len(df),
        "n_after": len(df_imputed),
        "rows_dropped": rows_dropped,
        "variables": valid_vars,
        "method": req.method.value,
        "missing_before": missing_before,
        "missing_after": missing_after,
        "imputation_details": imputation_details,
        "sample_changes": sample_changes[:10]
    }


@router.post("/imputation/apply")
async def apply_imputation(
    request: Request,
    req: ImputationRequest
):
    """Apply imputation and create a new snapshot with imputed data"""
    if not req.create_snapshot:
        raise HTTPException(status_code=400, detail="create_snapshot must be true to apply imputation")
    
    db = request.app.state.db
    
    # Get data
    if req.snapshot_id:
        snapshot_data = await db.snapshot_data.find_one({"snapshot_id": req.snapshot_id})
        if not snapshot_data:
            raise HTTPException(status_code=404, detail="Snapshot not found")
        data = snapshot_data.get("data", [])
        schema = snapshot_data.get("schema", [])
        source_name = f"Snapshot {req.snapshot_id[:8]}"
    else:
        submissions = await db.submissions.find({
            "form_id": req.form_id,
            "status": {"$in": ["approved", "submitted"]}
        }).to_list(None)
        data = [s.get("data", {}) for s in submissions]
        form = await db.forms.find_one({"id": req.form_id})
        schema = form.get("fields", []) if form else []
        source_name = f"Form {req.form_id[:8]}"
    
    if not data:
        return {"error": "No data available"}
    
    df = pd.DataFrame(data)
    valid_vars = [v for v in req.variables if v in df.columns]
    
    # Apply imputation (same logic as preview)
    for var in valid_vars:
        series = df[var].replace('', np.nan)
        try:
            numeric_series = pd.to_numeric(series, errors='coerce')
            is_numeric = numeric_series.notna().sum() > 0
        except Exception:
            is_numeric = False
            numeric_series = series
        
        if req.method == ImputationMethod.MEAN and is_numeric:
            df[var] = numeric_series.fillna(numeric_series.mean())
        elif req.method == ImputationMethod.MEDIAN and is_numeric:
            df[var] = numeric_series.fillna(numeric_series.median())
        elif req.method == ImputationMethod.MODE:
            mode_val = series.mode()
            if len(mode_val) > 0:
                df[var] = series.fillna(mode_val.iloc[0])
        elif req.method == ImputationMethod.CONSTANT:
            df[var] = series.fillna(req.constant_value)
        elif req.method == ImputationMethod.FORWARD_FILL:
            df[var] = series.ffill()
        elif req.method == ImputationMethod.BACKWARD_FILL:
            df[var] = series.bfill()
        elif req.method == ImputationMethod.INTERPOLATE and is_numeric:
            df[var] = numeric_series.interpolate(method='linear')
    
    if req.method == ImputationMethod.DROP:
        df = df.dropna(subset=valid_vars)
    
    # Create new snapshot
    import uuid
    snapshot_id = str(uuid.uuid4())
    
    snapshot = {
        "id": snapshot_id,
        "org_id": req.org_id,
        "form_id": req.form_id,
        "name": f"Imputed ({req.method.value}) - {source_name}",
        "description": f"Data with {req.method.value} imputation applied to: {', '.join(valid_vars)}",
        "status": "ready",
        "created_at": datetime.now(timezone.utc).isoformat(),
        "record_count": len(df),
        "transformation": {
            "type": "imputation",
            "method": req.method.value,
            "variables": valid_vars,
            "source_snapshot_id": req.snapshot_id
        }
    }
    
    snapshot_data_record = {
        "snapshot_id": snapshot_id,
        "schema": schema,
        "data": df.to_dict('records')
    }
    
    await db.snapshots.insert_one(snapshot)
    await db.snapshot_data.insert_one(snapshot_data_record)
    
    return {
        "snapshot_id": snapshot_id,
        "name": snapshot["name"],
        "record_count": len(df),
        "variables_imputed": valid_vars,
        "method": req.method.value
    }


@router.get("/imputation/missing-summary/{form_id}")
async def get_missing_summary(
    request: Request,
    form_id: str,
    snapshot_id: Optional[str] = None
):
    """Get summary of missing data across all variables"""
    db = request.app.state.db
    
    # Get data
    if snapshot_id:
        snapshot_data = await db.snapshot_data.find_one({"snapshot_id": snapshot_id})
        if not snapshot_data:
            raise HTTPException(status_code=404, detail="Snapshot not found")
        data = snapshot_data.get("data", [])
        schema = snapshot_data.get("schema", [])
    else:
        submissions = await db.submissions.find({
            "form_id": form_id,
            "status": {"$in": ["approved", "submitted"]}
        }).to_list(None)
        data = [s.get("data", {}) for s in submissions]
        form = await db.forms.find_one({"id": form_id})
        schema = form.get("fields", []) if form else []
    
    if not data:
        return {"total_rows": 0, "variables": [], "overall_missing_percent": 0}
    
    df = pd.DataFrame(data)
    
    variables_summary = []
    total_missing = 0
    total_cells = 0
    
    for col in df.columns:
        field_info = next((f for f in schema if f.get("id") == col), {})
        missing_count = df[col].isna().sum() + (df[col] == '').sum()
        total_missing += missing_count
        total_cells += len(df)
        
        # Determine if numeric
        try:
            pd.to_numeric(df[col], errors='raise')
            is_numeric = True
        except Exception:
            is_numeric = False
        
        variables_summary.append({
            "variable": col,
            "label": field_info.get("label", col),
            "type": field_info.get("type", "unknown"),
            "is_numeric": is_numeric,
            "total": len(df),
            "missing_count": int(missing_count),
            "missing_percent": round(float(missing_count / len(df) * 100), 2),
            "complete_count": int(len(df) - missing_count),
            "complete_percent": round(float((len(df) - missing_count) / len(df) * 100), 2)
        })
    
    # Sort by missing percent descending
    variables_summary.sort(key=lambda x: x["missing_percent"], reverse=True)
    
    return {
        "total_rows": len(df),
        "total_columns": len(df.columns),
        "total_cells": total_cells,
        "total_missing": int(total_missing),
        "overall_missing_percent": round(float(total_missing / total_cells * 100), 2) if total_cells > 0 else 0,
        "complete_cases": int((~df.isna().any(axis=1)).sum()),
        "complete_cases_percent": round(float((~df.isna().any(axis=1)).sum() / len(df) * 100), 2) if len(df) > 0 else 0,
        "variables": variables_summary
    }


# ============ Chart Data Endpoints ============

class CorrelationHeatmapRequest(BaseModel):
    org_id: str
    form_id: Optional[str] = None
    snapshot_id: Optional[str] = None
    variables: List[str]


@router.post("/charts/heatmap")
async def get_correlation_heatmap(
    request: Request,
    req: CorrelationHeatmapRequest
):
    """Get correlation matrix data for heatmap visualization"""
    db = request.app.state.db
    
    # Get data
    if req.snapshot_id:
        snapshot_data = await db.snapshot_data.find_one({"snapshot_id": req.snapshot_id})
        if not snapshot_data:
            raise HTTPException(status_code=404, detail="Snapshot not found")
        data = snapshot_data.get("data", [])
    else:
        submissions = await db.submissions.find({
            "form_id": req.form_id,
            "status": {"$in": ["approved", "submitted"]}
        }).to_list(None)
        data = [s.get("data", {}) for s in submissions]
    
    if not data:
        return {"error": "No data available"}
    
    df = pd.DataFrame(data)
    
    # Filter to requested variables and convert to numeric
    valid_vars = [v for v in req.variables if v in df.columns]
    if len(valid_vars) < 2:
        return {"error": "Need at least 2 variables for correlation heatmap"}
    
    df_numeric = df[valid_vars].apply(pd.to_numeric, errors='coerce').dropna()
    
    if len(df_numeric) < 5:
        return {"error": "Not enough data for correlation analysis"}
    
    # Calculate correlation matrix
    corr_matrix = df_numeric.corr()
    
    # Format for heatmap
    heatmap_data = []
    for i, var1 in enumerate(valid_vars):
        for j, var2 in enumerate(valid_vars):
            if var1 in corr_matrix.columns and var2 in corr_matrix.index:
                corr_val = corr_matrix.loc[var2, var1]
                heatmap_data.append({
                    "x": var1,
                    "y": var2,
                    "value": round(float(corr_val), 3) if not pd.isna(corr_val) else 0
                })
    
    return {
        "variables": valid_vars,
        "n_observations": len(df_numeric),
        "data": heatmap_data,
        "min_value": -1,
        "max_value": 1
    }


class ViolinPlotRequest(BaseModel):
    org_id: str
    form_id: Optional[str] = None
    snapshot_id: Optional[str] = None
    numeric_var: str
    group_var: Optional[str] = None


@router.post("/charts/violin")
async def get_violin_data(
    request: Request,
    req: ViolinPlotRequest
):
    """Get data for violin plot visualization (distribution by group)"""
    db = request.app.state.db
    
    # Get data
    if req.snapshot_id:
        snapshot_data = await db.snapshot_data.find_one({"snapshot_id": req.snapshot_id})
        if not snapshot_data:
            raise HTTPException(status_code=404, detail="Snapshot not found")
        data = snapshot_data.get("data", [])
    else:
        submissions = await db.submissions.find({
            "form_id": req.form_id,
            "status": {"$in": ["approved", "submitted"]}
        }).to_list(None)
        data = [s.get("data", {}) for s in submissions]
    
    if not data:
        return {"error": "No data available"}
    
    df = pd.DataFrame(data)
    
    if req.numeric_var not in df.columns:
        return {"error": f"Variable {req.numeric_var} not found"}
    
    # Convert to numeric
    df[req.numeric_var] = pd.to_numeric(df[req.numeric_var], errors='coerce')
    df = df.dropna(subset=[req.numeric_var])
    
    if len(df) < 10:
        return {"error": "Not enough data for violin plot"}
    
    # Calculate distribution data
    result = {
        "variable": req.numeric_var,
        "n_total": len(df),
        "groups": []
    }
    
    if req.group_var and req.group_var in df.columns:
        # Group-wise statistics
        for group_name, group_df in df.groupby(req.group_var):
            values = group_df[req.numeric_var].dropna()
            if len(values) >= 5:
                # Calculate percentiles for violin shape
                percentiles = [5, 10, 25, 50, 75, 90, 95]
                pcts = np.percentile(values, percentiles)
                
                result["groups"].append({
                    "name": str(group_name),
                    "n": len(values),
                    "mean": round(float(values.mean()), 3),
                    "median": round(float(values.median()), 3),
                    "std": round(float(values.std()), 3),
                    "min": round(float(values.min()), 3),
                    "max": round(float(values.max()), 3),
                    "q1": round(float(pcts[2]), 3),  # 25th percentile
                    "q3": round(float(pcts[4]), 3),  # 75th percentile
                    "p5": round(float(pcts[0]), 3),
                    "p95": round(float(pcts[6]), 3),
                    "distribution": [
                        {"percentile": p, "value": round(float(v), 3)}
                        for p, v in zip(percentiles, pcts)
                    ]
                })
    else:
        # Overall statistics
        values = df[req.numeric_var].dropna()
        percentiles = [5, 10, 25, 50, 75, 90, 95]
        pcts = np.percentile(values, percentiles)
        
        result["groups"].append({
            "name": "All",
            "n": len(values),
            "mean": round(float(values.mean()), 3),
            "median": round(float(values.median()), 3),
            "std": round(float(values.std()), 3),
            "min": round(float(values.min()), 3),
            "max": round(float(values.max()), 3),
            "q1": round(float(pcts[2]), 3),
            "q3": round(float(pcts[4]), 3),
            "p5": round(float(pcts[0]), 3),
            "p95": round(float(pcts[6]), 3),
            "distribution": [
                {"percentile": p, "value": round(float(v), 3)}
                for p, v in zip(percentiles, pcts)
            ]
        })
    
    return result


class CoefficientPlotRequest(BaseModel):
    org_id: str
    form_id: Optional[str] = None
    snapshot_id: Optional[str] = None
    dependent_var: str
    independent_vars: List[str]


@router.post("/charts/coefficient")
async def get_coefficient_plot_data(
    request: Request,
    req: CoefficientPlotRequest
):
    """Get regression coefficients with confidence intervals for coefficient plot"""
    db = request.app.state.db
    
    # Get data
    if req.snapshot_id:
        snapshot_data = await db.snapshot_data.find_one({"snapshot_id": req.snapshot_id})
        if not snapshot_data:
            raise HTTPException(status_code=404, detail="Snapshot not found")
        data = snapshot_data.get("data", [])
    else:
        submissions = await db.submissions.find({
            "form_id": req.form_id,
            "status": {"$in": ["approved", "submitted"]}
        }).to_list(None)
        data = [s.get("data", {}) for s in submissions]
    
    if not data:
        return {"error": "No data available"}
    
    df = pd.DataFrame(data)
    
    # Prepare variables
    all_vars = [req.dependent_var] + req.independent_vars
    valid_vars = [v for v in all_vars if v in df.columns]
    
    if len(valid_vars) < 2:
        return {"error": "Need at least dependent variable and one independent variable"}
    
    # Convert to numeric and drop missing
    df_model = df[valid_vars].apply(pd.to_numeric, errors='coerce').dropna()
    
    if len(df_model) < len(req.independent_vars) + 10:
        return {"error": "Not enough data for regression"}
    
    # Run OLS regression
    import statsmodels.api as sm
    
    y = df_model[req.dependent_var]
    X = df_model[[v for v in req.independent_vars if v in df_model.columns]]
    X = sm.add_constant(X)
    
    try:
        model = sm.OLS(y, X).fit()
        
        coefficients = []
        for var in model.params.index:
            if var != 'const':
                coef = model.params[var]
                se = model.bse[var]
                conf_int = model.conf_int().loc[var]
                p_val = model.pvalues[var]
                
                coefficients.append({
                    "variable": var,
                    "coefficient": round(float(coef), 4),
                    "std_error": round(float(se), 4),
                    "ci_lower": round(float(conf_int[0]), 4),
                    "ci_upper": round(float(conf_int[1]), 4),
                    "p_value": round(float(p_val), 4),
                    "significant": p_val < 0.05
                })
        
        # Sort by absolute coefficient size
        coefficients.sort(key=lambda x: abs(x["coefficient"]), reverse=True)
        
        return {
            "dependent_var": req.dependent_var,
            "n_observations": len(df_model),
            "r_squared": round(float(model.rsquared), 4),
            "adj_r_squared": round(float(model.rsquared_adj), 4),
            "coefficients": coefficients,
            "intercept": {
                "value": round(float(model.params.get('const', 0)), 4),
                "ci_lower": round(float(model.conf_int().loc['const'][0]), 4) if 'const' in model.conf_int().index else None,
                "ci_upper": round(float(model.conf_int().loc['const'][1]), 4) if 'const' in model.conf_int().index else None
            }
        }
    except Exception as e:
        return {"error": f"Regression failed: {str(e)}"}



# ============ Chart Export ============

class ChartExportRequest(BaseModel):
    image_data: str  # Base64 PNG data
    title: str
    subtitle: Optional[str] = None
    width: int = 800
    height: int = 600


@router.post("/export-chart-pdf")
async def export_chart_pdf(request: Request, req: ChartExportRequest):
    """Export chart as PDF with title and subtitle"""
    from fastapi.responses import StreamingResponse
    import base64
    
    try:
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.pdfgen import canvas
        from reportlab.lib.units import inch
        from reportlab.lib import colors
    except ImportError:
        raise HTTPException(status_code=500, detail="PDF generation not available")
    
    # Decode base64 image
    try:
        image_data = req.image_data.split(',')[1] if ',' in req.image_data else req.image_data
        image_bytes = base64.b64decode(image_data)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid image data")
    
    # Create PDF
    buffer = io.BytesIO()
    
    # Use landscape if chart is wider than tall
    if req.width > req.height:
        page_size = landscape(letter)
    else:
        page_size = letter
    
    c = canvas.Canvas(buffer, pagesize=page_size)
    page_width, page_height = page_size
    
    # Title
    c.setFont("Helvetica-Bold", 18)
    c.drawCentredString(page_width / 2, page_height - 0.75 * inch, req.title)
    
    # Subtitle
    if req.subtitle:
        c.setFont("Helvetica", 12)
        c.setFillColor(colors.gray)
        c.drawCentredString(page_width / 2, page_height - 1.1 * inch, req.subtitle)
        c.setFillColor(colors.black)
    
    # Calculate image size to fit page with margins
    margin = 0.75 * inch
    title_space = 1.5 * inch if req.subtitle else 1.2 * inch
    available_width = page_width - 2 * margin
    available_height = page_height - title_space - margin
    
    # Maintain aspect ratio
    aspect = req.width / req.height
    if available_width / aspect <= available_height:
        img_width = available_width
        img_height = available_width / aspect
    else:
        img_height = available_height
        img_width = available_height * aspect
    
    # Center the image
    x = (page_width - img_width) / 2
    y = page_height - title_space - img_height
    
    # Draw image from bytes
    from reportlab.lib.utils import ImageReader
    from PIL import Image
    
    img = Image.open(io.BytesIO(image_bytes))
    img_reader = ImageReader(img)
    c.drawImage(img_reader, x, y, width=img_width, height=img_height)
    
    # Footer with timestamp
    c.setFont("Helvetica", 8)
    c.setFillColor(colors.gray)
    timestamp = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    c.drawString(margin, 0.4 * inch, f"Generated: {timestamp}")
    c.drawRightString(page_width - margin, 0.4 * inch, "DataPulse Analytics")
    
    c.save()
    buffer.seek(0)
    
    return StreamingResponse(
        buffer,
        media_type="application/pdf",
        headers={
            "Content-Disposition": f"attachment; filename=chart_{req.title.replace(' ', '_')}.pdf"
        }
    )


@router.post("/export-chart-pptx")
async def export_chart_pptx(request: Request, req: ChartExportRequest):
    """Export chart as PowerPoint slide"""
    import base64
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx required for PowerPoint export. Install with: pip install python-pptx")
    
    # Decode base64 image
    try:
        image_data = req.image_data.split(',')[1] if ',' in req.image_data else req.image_data
        image_bytes = base64.b64decode(image_data)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid image data")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)
    
    # Add blank slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = req.title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle if present
    if req.subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12.333), Inches(0.4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = req.subtitle
        subtitle_para.font.size = Pt(16)
        subtitle_para.font.color.rgb = RgbColor(100, 100, 100)
        subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Save image to temp file and add to slide
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(image_bytes)
        tmp_path = tmp.name
    
    # Calculate image position to center it
    img_top = Inches(1.5) if req.subtitle else Inches(1.2)
    available_height = Inches(5.5) if req.subtitle else Inches(5.8)
    available_width = Inches(12.333)
    
    # Maintain aspect ratio
    aspect = req.width / req.height
    if available_width / aspect <= available_height:
        img_width = available_width
        img_height = available_width / aspect
    else:
        img_height = available_height
        img_width = available_height * aspect
    
    img_left = (Inches(13.333) - img_width) / 2
    
    slide.shapes.add_picture(tmp_path, img_left, img_top, width=img_width, height=img_height)
    
    # Add footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.333), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_para = footer_frame.paragraphs[0]
    timestamp = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    footer_para.text = f"DataPulse Analytics | Generated: {timestamp}"
    footer_para.font.size = Pt(10)
    footer_para.font.color.rgb = RgbColor(150, 150, 150)
    footer_para.alignment = PP_ALIGN.CENTER
    
    # Clean up temp file
    import os
    os.unlink(tmp_path)
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename=chart_{req.title.replace(' ', '_')}.pptx"
        }
    )


# ============ AI-Assisted Data Preparation ============

class DataPrepRequest(BaseModel):
    snapshot_id: Optional[str] = None
    form_id: Optional[str] = None
    org_id: str
    variables: Optional[List[str]] = None  # If None, analyze all variables


class SuggestionPriority(str, Enum):
    HIGH = "high"
    MEDIUM = "medium"
    LOW = "low"


@router.post("/data-prep-suggestions")
async def get_data_prep_suggestions(request: Request, req: DataPrepRequest):
    """
    AI-Assisted Data Preparation: Analyze dataset and provide intelligent suggestions
    for data cleaning, transformations, and recoding.
    """
    db = request.app.state.db
    
    # Get data
    if req.snapshot_id:
        snapshot = await db.snapshots.find_one({"id": req.snapshot_id}, {"_id": 0})
        if not snapshot:
            raise HTTPException(status_code=404, detail="Snapshot not found")
        df = pd.DataFrame(snapshot.get("data", []))
        schema = snapshot.get("schema", {})
    elif req.form_id:
        form = await db.forms.find_one({"id": req.form_id}, {"_id": 0})
        if not form:
            raise HTTPException(status_code=404, detail="Form not found")
        
        responses = await db.responses.find(
            {"form_id": req.form_id, "is_deleted": {"$ne": True}},
            {"_id": 0}
        ).to_list(10000)
        
        if not responses:
            return {"suggestions": [], "summary": "No data available for analysis"}
        
        df = pd.DataFrame([r.get("data", {}) for r in responses])
        schema = {f["id"]: f for f in form.get("fields", [])}
    else:
        raise HTTPException(status_code=400, detail="Either snapshot_id or form_id required")
    
    if df.empty:
        return {"suggestions": [], "summary": "No data available for analysis"}
    
    # Filter variables if specified
    if req.variables:
        available = [v for v in req.variables if v in df.columns]
        if not available:
            raise HTTPException(status_code=400, detail="No specified variables found in data")
        df = df[available]
    
    suggestions = []
    
    # Analyze each variable
    for col in df.columns:
        col_suggestions = analyze_variable(df[col], col, schema.get(col, {}))
        suggestions.extend(col_suggestions)
    
    # Dataset-level suggestions
    dataset_suggestions = analyze_dataset(df, schema)
    suggestions.extend(dataset_suggestions)
    
    # Sort by priority
    priority_order = {"high": 0, "medium": 1, "low": 2}
    suggestions.sort(key=lambda x: priority_order.get(x.get("priority", "low"), 2))
    
    # Generate summary
    high_count = len([s for s in suggestions if s.get("priority") == "high"])
    medium_count = len([s for s in suggestions if s.get("priority") == "medium"])
    
    summary = f"Found {len(suggestions)} suggestions: {high_count} high priority, {medium_count} medium priority."
    if high_count > 0:
        summary += " Address high-priority items first for data quality."
    
    return {
        "suggestions": suggestions,
        "summary": summary,
        "variables_analyzed": len(df.columns),
        "observations": len(df)
    }


def analyze_variable(series: pd.Series, var_name: str, schema_info: dict) -> List[dict]:
    """Analyze a single variable and return suggestions"""
    suggestions = []
    var_label = schema_info.get("label", var_name)
    var_type = schema_info.get("type", "unknown")
    
    # 1. Missing Data Analysis
    missing_count = series.isna().sum()
    missing_pct = (missing_count / len(series)) * 100 if len(series) > 0 else 0
    
    if missing_pct > 50:
        suggestions.append({
            "type": "missing_data",
            "priority": "high",
            "variable": var_name,
            "variable_label": var_label,
            "title": f"High missing data in '{var_label}'",
            "description": f"{missing_pct:.1f}% of values are missing ({missing_count} of {len(series)}). Consider removing this variable or using imputation.",
            "actions": [
                {"action": "impute_mean", "label": "Impute with mean"},
                {"action": "impute_median", "label": "Impute with median"},
                {"action": "impute_mode", "label": "Impute with mode"},
                {"action": "drop_variable", "label": "Drop variable"}
            ]
        })
    elif missing_pct > 10:
        suggestions.append({
            "type": "missing_data",
            "priority": "medium",
            "variable": var_name,
            "variable_label": var_label,
            "title": f"Moderate missing data in '{var_label}'",
            "description": f"{missing_pct:.1f}% of values are missing. Consider imputation or listwise deletion.",
            "actions": [
                {"action": "impute_mean", "label": "Impute with mean"},
                {"action": "impute_median", "label": "Impute with median"},
                {"action": "drop_missing", "label": "Drop rows with missing"}
            ]
        })
    elif missing_pct > 0:
        suggestions.append({
            "type": "missing_data",
            "priority": "low",
            "variable": var_name,
            "variable_label": var_label,
            "title": f"Some missing data in '{var_label}'",
            "description": f"{missing_pct:.1f}% missing values ({missing_count}). May be acceptable for most analyses.",
            "actions": [
                {"action": "drop_missing", "label": "Drop rows with missing"},
                {"action": "ignore", "label": "Ignore (use listwise deletion)"}
            ]
        })
    
    # Work with non-missing values
    valid_series = series.dropna()
    if len(valid_series) == 0:
        return suggestions
    
    # 2. Numeric Variable Analysis
    numeric_series = pd.to_numeric(valid_series, errors='coerce')
    if numeric_series.notna().sum() > len(valid_series) * 0.8:  # Mostly numeric
        numeric_valid = numeric_series.dropna()
        
        if len(numeric_valid) > 10:
            # Outlier detection using IQR
            q1 = numeric_valid.quantile(0.25)
            q3 = numeric_valid.quantile(0.75)
            iqr = q3 - q1
            lower_bound = q1 - 1.5 * iqr
            upper_bound = q3 + 1.5 * iqr
            
            outliers = numeric_valid[(numeric_valid < lower_bound) | (numeric_valid > upper_bound)]
            outlier_pct = (len(outliers) / len(numeric_valid)) * 100
            
            if outlier_pct > 10:
                suggestions.append({
                    "type": "outliers",
                    "priority": "high",
                    "variable": var_name,
                    "variable_label": var_label,
                    "title": f"Many outliers in '{var_label}'",
                    "description": f"{outlier_pct:.1f}% of values are outliers ({len(outliers)} values outside [{lower_bound:.2f}, {upper_bound:.2f}]). This may affect statistical analyses.",
                    "stats": {
                        "outlier_count": int(len(outliers)),
                        "lower_bound": round(float(lower_bound), 2),
                        "upper_bound": round(float(upper_bound), 2),
                        "min_outlier": round(float(outliers.min()), 2) if len(outliers) > 0 else None,
                        "max_outlier": round(float(outliers.max()), 2) if len(outliers) > 0 else None
                    },
                    "actions": [
                        {"action": "winsorize", "label": "Winsorize (cap at bounds)"},
                        {"action": "remove_outliers", "label": "Remove outliers"},
                        {"action": "log_transform", "label": "Log transform"}
                    ]
                })
            elif outlier_pct > 2:
                suggestions.append({
                    "type": "outliers",
                    "priority": "low",
                    "variable": var_name,
                    "variable_label": var_label,
                    "title": f"Some outliers in '{var_label}'",
                    "description": f"{len(outliers)} outlier(s) detected. Review if they are valid data points.",
                    "stats": {
                        "outlier_count": int(len(outliers)),
                        "values": [round(float(v), 2) for v in outliers.head(5).tolist()]
                    },
                    "actions": [
                        {"action": "review", "label": "Review manually"},
                        {"action": "winsorize", "label": "Winsorize"}
                    ]
                })
            
            # Skewness check
            from scipy import stats as scipy_stats
            skewness = scipy_stats.skew(numeric_valid)
            
            if abs(skewness) > 2:
                suggestions.append({
                    "type": "distribution",
                    "priority": "medium",
                    "variable": var_name,
                    "variable_label": var_label,
                    "title": f"Highly skewed distribution in '{var_label}'",
                    "description": f"Skewness = {skewness:.2f}. {'Positively' if skewness > 0 else 'Negatively'} skewed data may violate normality assumptions.",
                    "stats": {
                        "skewness": round(float(skewness), 2),
                        "mean": round(float(numeric_valid.mean()), 2),
                        "median": round(float(numeric_valid.median()), 2)
                    },
                    "actions": [
                        {"action": "log_transform", "label": "Log transform"} if numeric_valid.min() > 0 else None,
                        {"action": "sqrt_transform", "label": "Square root transform"} if numeric_valid.min() >= 0 else None,
                        {"action": "box_cox", "label": "Box-Cox transform"} if numeric_valid.min() > 0 else None,
                        {"action": "ignore", "label": "Use non-parametric tests"}
                    ]
                })
            
            # Zero variance check
            if numeric_valid.std() == 0:
                suggestions.append({
                    "type": "zero_variance",
                    "priority": "high",
                    "variable": var_name,
                    "variable_label": var_label,
                    "title": f"Zero variance in '{var_label}'",
                    "description": f"All values are identical ({numeric_valid.iloc[0]}). This variable provides no information.",
                    "actions": [
                        {"action": "drop_variable", "label": "Drop variable"}
                    ]
                })
    
    # 3. Categorical Variable Analysis
    else:
        unique_values = valid_series.unique()
        value_counts = valid_series.value_counts()
        
        # Check for inconsistent labels (case variations, whitespace)
        str_values = [str(v).strip() for v in unique_values]
        lower_values = [v.lower() for v in str_values]
        
        # Find potential duplicates due to case/whitespace
        seen = {}
        inconsistent = []
        for orig, lower in zip(str_values, lower_values):
            if lower in seen and seen[lower] != orig:
                inconsistent.append((seen[lower], orig))
            seen[lower] = orig
        
        if inconsistent:
            suggestions.append({
                "type": "inconsistent_labels",
                "priority": "high",
                "variable": var_name,
                "variable_label": var_label,
                "title": f"Inconsistent labels in '{var_label}'",
                "description": f"Found similar values that may be the same category: {inconsistent[:3]}",
                "examples": inconsistent[:5],
                "actions": [
                    {"action": "standardize_case", "label": "Standardize to lowercase"},
                    {"action": "trim_whitespace", "label": "Trim whitespace"},
                    {"action": "review_recode", "label": "Review and recode manually"}
                ]
            })
        
        # Check for rare categories
        total = len(valid_series)
        rare_categories = value_counts[value_counts < total * 0.02]  # Less than 2%
        
        if len(rare_categories) > 0 and len(rare_categories) < len(value_counts):
            suggestions.append({
                "type": "rare_categories",
                "priority": "low",
                "variable": var_name,
                "variable_label": var_label,
                "title": f"Rare categories in '{var_label}'",
                "description": f"{len(rare_categories)} categories have less than 2% of responses. Consider combining into 'Other'.",
                "categories": {str(k): int(v) for k, v in rare_categories.head(5).items()},
                "actions": [
                    {"action": "combine_rare", "label": "Combine into 'Other'"},
                    {"action": "ignore", "label": "Keep as is"}
                ]
            })
        
        # Too many unique values for categorical
        if len(unique_values) > 20 and var_type in ['select', 'radio', 'categorical']:
            suggestions.append({
                "type": "high_cardinality",
                "priority": "medium",
                "variable": var_name,
                "variable_label": var_label,
                "title": f"High cardinality in '{var_label}'",
                "description": f"{len(unique_values)} unique values. May need grouping for meaningful analysis.",
                "actions": [
                    {"action": "group_values", "label": "Create value groups"},
                    {"action": "review", "label": "Review manually"}
                ]
            })
    
    # Filter out None actions
    for s in suggestions:
        if "actions" in s:
            s["actions"] = [a for a in s["actions"] if a is not None]
    
    return suggestions


def analyze_dataset(df: pd.DataFrame, schema: dict) -> List[dict]:
    """Analyze dataset-level issues"""
    suggestions = []
    
    # 1. Sample size check
    if len(df) < 30:
        suggestions.append({
            "type": "sample_size",
            "priority": "high",
            "variable": None,
            "variable_label": None,
            "title": "Small sample size",
            "description": f"Only {len(df)} observations. Many statistical tests require n ≥ 30 for reliable results.",
            "actions": [
                {"action": "use_nonparametric", "label": "Use non-parametric tests"},
                {"action": "acknowledge", "label": "Acknowledge limitation"}
            ]
        })
    
    # 2. Check for duplicate rows
    duplicates = df.duplicated().sum()
    if duplicates > 0:
        dup_pct = (duplicates / len(df)) * 100
        suggestions.append({
            "type": "duplicates",
            "priority": "medium" if dup_pct > 5 else "low",
            "variable": None,
            "variable_label": None,
            "title": "Duplicate rows detected",
            "description": f"{duplicates} duplicate rows ({dup_pct:.1f}%). May indicate data collection issues.",
            "actions": [
                {"action": "remove_duplicates", "label": "Remove duplicates"},
                {"action": "review", "label": "Review duplicates"}
            ]
        })
    
    # 3. Check for highly correlated numeric variables
    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    if len(numeric_cols) >= 2:
        try:
            corr_matrix = df[numeric_cols].corr()
            high_corr = []
            
            for i, col1 in enumerate(numeric_cols):
                for col2 in numeric_cols[i+1:]:
                    corr_val = corr_matrix.loc[col1, col2]
                    if abs(corr_val) > 0.9 and not pd.isna(corr_val):
                        high_corr.append({
                            "var1": col1,
                            "var2": col2,
                            "correlation": round(float(corr_val), 3)
                        })
            
            if high_corr:
                suggestions.append({
                    "type": "multicollinearity",
                    "priority": "medium",
                    "variable": None,
                    "variable_label": None,
                    "title": "Highly correlated variables",
                    "description": f"{len(high_corr)} variable pair(s) with correlation > 0.9. May cause multicollinearity in regression.",
                    "pairs": high_corr[:5],
                    "actions": [
                        {"action": "remove_one", "label": "Remove one variable from each pair"},
                        {"action": "pca", "label": "Use PCA to combine"},
                        {"action": "acknowledge", "label": "Acknowledge in analysis"}
                    ]
                })
        except Exception:
            pass
    
    # 4. ID-like columns that shouldn't be analyzed
    for col in df.columns:
        if len(df[col].dropna()) > 0:
            unique_ratio = df[col].nunique() / len(df[col].dropna())
            col_lower = col.lower()
            
            if unique_ratio > 0.95 and any(kw in col_lower for kw in ['id', 'uuid', 'key', 'index', 'email', 'phone']):
                suggestions.append({
                    "type": "identifier_column",
                    "priority": "low",
                    "variable": col,
                    "variable_label": schema.get(col, {}).get("label", col),
                    "title": f"Possible identifier column: '{col}'",
                    "description": "This column appears to be an identifier (nearly all unique values). Exclude from statistical analysis.",
                    "actions": [
                        {"action": "exclude", "label": "Exclude from analysis"},
                        {"action": "ignore", "label": "Keep in dataset"}
                    ]
                })
    
    return suggestions


@router.post("/apply-transformation")
async def apply_transformation(request: Request, transformation: dict):
    """Apply a data transformation based on suggestion action"""
    db = request.app.state.db
    
    action = transformation.get("action")
    variable = transformation.get("variable")
    snapshot_id = transformation.get("snapshot_id")
    
    if not action:
        raise HTTPException(status_code=400, detail="Action required")
    
    # Get data
    if snapshot_id:
        snapshot = await db.snapshots.find_one({"id": snapshot_id})
        if not snapshot:
            raise HTTPException(status_code=404, detail="Snapshot not found")
        df = pd.DataFrame(snapshot.get("data", []))
    else:
        raise HTTPException(status_code=400, detail="Snapshot ID required for transformations")
    
    if df.empty:
        raise HTTPException(status_code=400, detail="No data to transform")
    
    # Apply transformation
    result = {"success": False, "message": "", "affected_rows": 0}
    
    try:
        if action == "impute_mean" and variable:
            if variable in df.columns:
                mean_val = pd.to_numeric(df[variable], errors='coerce').mean()
                affected = df[variable].isna().sum()
                df[variable] = pd.to_numeric(df[variable], errors='coerce').fillna(mean_val)
                result = {"success": True, "message": f"Imputed {affected} values with mean ({mean_val:.2f})", "affected_rows": int(affected)}
        
        elif action == "impute_median" and variable:
            if variable in df.columns:
                median_val = pd.to_numeric(df[variable], errors='coerce').median()
                affected = df[variable].isna().sum()
                df[variable] = pd.to_numeric(df[variable], errors='coerce').fillna(median_val)
                result = {"success": True, "message": f"Imputed {affected} values with median ({median_val:.2f})", "affected_rows": int(affected)}
        
        elif action == "impute_mode" and variable:
            if variable in df.columns:
                mode_val = df[variable].mode().iloc[0] if len(df[variable].mode()) > 0 else None
                if mode_val is not None:
                    affected = df[variable].isna().sum()
                    df[variable] = df[variable].fillna(mode_val)
                    result = {"success": True, "message": f"Imputed {affected} values with mode ({mode_val})", "affected_rows": int(affected)}
        
        elif action == "drop_missing" and variable:
            if variable in df.columns:
                original_len = len(df)
                df = df.dropna(subset=[variable])
                affected = original_len - len(df)
                result = {"success": True, "message": f"Removed {affected} rows with missing values", "affected_rows": affected}
        
        elif action == "remove_duplicates":
            original_len = len(df)
            df = df.drop_duplicates()
            affected = original_len - len(df)
            result = {"success": True, "message": f"Removed {affected} duplicate rows", "affected_rows": affected}
        
        elif action == "standardize_case" and variable:
            if variable in df.columns:
                df[variable] = df[variable].astype(str).str.lower().str.strip()
                result = {"success": True, "message": f"Standardized case for '{variable}'", "affected_rows": len(df)}
        
        elif action == "trim_whitespace" and variable:
            if variable in df.columns:
                df[variable] = df[variable].astype(str).str.strip()
                result = {"success": True, "message": f"Trimmed whitespace for '{variable}'", "affected_rows": len(df)}
        
        elif action == "log_transform" and variable:
            if variable in df.columns:
                numeric_col = pd.to_numeric(df[variable], errors='coerce')
                if numeric_col.min() > 0:
                    df[f"{variable}_log"] = np.log(numeric_col)
                    result = {"success": True, "message": f"Created log-transformed variable '{variable}_log'", "affected_rows": len(df)}
                else:
                    result = {"success": False, "message": "Cannot log transform: contains zero or negative values"}
        
        elif action == "winsorize" and variable:
            if variable in df.columns:
                numeric_col = pd.to_numeric(df[variable], errors='coerce')
                q1 = numeric_col.quantile(0.25)
                q3 = numeric_col.quantile(0.75)
                iqr = q3 - q1
                lower = q1 - 1.5 * iqr
                upper = q3 + 1.5 * iqr
                
                affected = ((numeric_col < lower) | (numeric_col > upper)).sum()
                df[variable] = numeric_col.clip(lower=lower, upper=upper)
                result = {"success": True, "message": f"Winsorized {affected} outliers to [{lower:.2f}, {upper:.2f}]", "affected_rows": int(affected)}
        
        else:
            result = {"success": False, "message": f"Unknown action: {action}"}
        
        # Save transformed data back to snapshot
        if result["success"]:
            await db.snapshots.update_one(
                {"id": snapshot_id},
                {
                    "$set": {
                        "data": df.to_dict(orient="records"),
                        "row_count": len(df),
                        "updated_at": datetime.now(timezone.utc)
                    },
                    "$push": {
                        "transformations": {
                            "action": action,
                            "variable": variable,
                            "timestamp": datetime.now(timezone.utc).isoformat(),
                            "affected_rows": result["affected_rows"]
                        }
                    }
                }
            )
        
        return result
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Transformation failed: {str(e)}")

